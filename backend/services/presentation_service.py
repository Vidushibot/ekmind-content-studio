import json
import re
import shutil
import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from backend.schemas.content import Slide


def _humanized_fallback_notes(title: str, points: list[str], slide_number: int, total: int) -> str:
    ideas = [re.sub(r"\[([^]]+)\]\([^)]+\)", r"\1", point).strip().rstrip(".") for point in points if point.strip()]
    main = ideas[0] if ideas else title
    support = ideas[1] if len(ideas) > 1 else "the practical implications behind it"
    if slide_number == 1:
        return f"Here’s the idea I want to explore with you: {main}. It sounds straightforward, but the real value is in what it changes for people doing the work. We’ll move from the familiar assumption to a more useful, practical way of thinking about it."
    if slide_number == total:
        return f"So this is the takeaway I’d leave with you: {main}. Keep {support.lower()} in mind, and choose one realistic place to apply it. A small, deliberate next step will tell us far more than another broad promise."
    return f"The important point here is {main.lower()}. In practice, that means paying attention to {support.lower()}. It’s less about adding complexity and more about making the next decision clearer. That brings us naturally to the next part of the story."


def _clean_speaker_notes(text: str) -> str:
    text = re.sub(r"\[([^]]+)\]\([^)]+\)", r"\1", text or "")
    text = re.sub(r"https?://\S+", "", text)
    return " ".join(text.split()).strip()


def build_outline(topic: str, post: str, design_instructions: str = "") -> list[Slide]:
    paragraphs = [p.strip() for p in post.split("\n\n") if p.strip()]
    titles = [topic, "The tension", "A better question", "A practical path", "Human accountability", "Start bounded", "What changes", "Your next move"]
    types = ["TITLE", "TITLE_BODY", "KEY_MESSAGE", "PROCESS", "TWO_COLUMN", "KEY_MESSAGE", "QUOTE", "CLOSING"]
    slides = []
    for i, title in enumerate(titles):
        source = paragraphs[min(i, len(paragraphs) - 1)]
        clean_source = re.sub(r"\[([^]]+)\]\([^)]+\)", r"\1", source)
        candidates = [part.strip(" •0123456789.-") for part in re.split(r"(?<=[.!?])\s+|\n+", clean_source) if part.strip()]
        points = []
        for candidate in candidates[:3]:
            if len(candidate) <= 88:
                points.append(candidate)
        if not points:
            points = ["See speaker notes for the complete narrative."]
        visual_types = ["concept", "comparison", "quote", "process", "concept", "timeline", "metric", "closing"]
        notes = _humanized_fallback_notes(title, points, i + 1, len(titles))
        slides.append(Slide(slide_number=i + 1, slide_type=types[i], title=title, content=points, speaker_notes=notes, visual_type=visual_types[i], visual_labels=points[:3]))
    return slides


# Design tokens derived from the supplied HR Transformation Mythbuster deck.
NAVY = RGBColor(23, 40, 74)
INK = RGBColor(12, 24, 48)
GOLD = RGBColor(244, 178, 26)
CORAL = RGBColor(228, 87, 61)
PALE = RGBColor(255, 247, 224)
IVORY = RGBColor(250, 248, 242)
WHITE = RGBColor(255, 255, 255)


def _label(shape, text: str, size: int = 16, color: RGBColor = NAVY) -> None:
    shape.text = text
    frame = shape.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for paragraph in frame.paragraphs:
        paragraph.font.name = "Arial Black"
        paragraph.font.size = Pt(size)
        paragraph.font.bold = True
        paragraph.font.color.rgb = color
        for run in paragraph.runs:
            run.font.name = "Arial Black"
            run.font.size = Pt(size)
            run.font.bold = True
            run.font.color.rgb = color


def _textbox(slide, text: str, left: float, top: float, width: float, height: float, size: int, color: RGBColor, bold: bool = False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.text = text
    frame = box.text_frame; frame.word_wrap = True; frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = frame.margin_right = Inches(0)
    for paragraph in frame.paragraphs:
        paragraph.font.name = "Arial Black" if bold else "Arial"
        paragraph.font.size = Pt(size); paragraph.font.bold = bold; paragraph.font.color.rgb = color
        for run in paragraph.runs:
            run.font.name = "Arial Black" if bold else "Arial"
            run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return box


def _background(slide, color: RGBColor) -> None:
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color


def _add_visual(slide, item: Slide) -> None:
    labels = (item.visual_labels or item.content or [item.title])[:3]
    visual = item.visual_type
    if visual in {"process", "timeline"}:
        for index, label in enumerate(labels):
            x = 5.65 + index * 1.42
            node = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(2.08), Inches(1.34), Inches(1.18))
            node.fill.solid(); node.fill.fore_color.rgb = GOLD if index == 0 else NAVY; node.line.fill.background()
            _label(node, f"{index + 1}\n{label}", 11, INK if index == 0 else WHITE)
    elif visual == "comparison":
        for index, label in enumerate(labels[:2]):
            panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.72 + index * 1.95), Inches(1.82), Inches(1.7), Inches(2.55))
            panel.fill.solid(); panel.fill.fore_color.rgb = NAVY if index == 0 else GOLD; panel.line.fill.background()
            _label(panel, label, 14, WHITE if index == 0 else INK)
    elif visual == "quote":
        mark = slide.shapes.add_textbox(Inches(5.65), Inches(1.25), Inches(1.2), Inches(1.3)); _label(mark, "“", 66, GOLD)
        caption = slide.shapes.add_textbox(Inches(6.25), Inches(2.05), Inches(3.05), Inches(1.65)); _label(caption, labels[0], 17, INK)
    elif visual == "metric":
        metric = next((token for token in " ".join(labels).split() if any(char.isdigit() for char in token)), "ONE")
        halo = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.12), Inches(1.5), Inches(2.65), Inches(2.65)); halo.fill.solid(); halo.fill.fore_color.rgb = GOLD; halo.line.fill.background()
        number = slide.shapes.add_textbox(Inches(6.38), Inches(2.03), Inches(2.15), Inches(0.95)); _label(number, metric, 38, INK)
        caption = slide.shapes.add_textbox(Inches(5.75), Inches(3.72), Inches(3.4), Inches(0.75)); _label(caption, labels[0], 13, INK)
    elif visual == "closing":
        target = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.15), Inches(1.5), Inches(2.55), Inches(2.55))
        target.fill.solid(); target.fill.fore_color.rgb = PALE; target.line.color.rgb = GOLD; target.line.width = Pt(4)
        inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.82), Inches(2.17), Inches(1.2), Inches(1.2))
        inner.fill.solid(); inner.fill.fore_color.rgb = GOLD; inner.line.fill.background(); _label(inner, "→", 26, INK)
    else:
        positions = [(5.55, 1.4), (8.1, 1.4), (6.82, 3.92)]
        center = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.78), Inches(2.18), Inches(1.55), Inches(1.55))
        center.fill.solid(); center.fill.fore_color.rgb = NAVY; center.line.color.rgb = GOLD; center.line.width = Pt(3); _label(center, "CORE\nIDEA", 16, GOLD)
        for label, (x, y) in zip(labels, positions):
            node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(1.45), Inches(0.76))
            node.fill.solid(); node.fill.fore_color.rgb = GOLD; node.line.fill.background(); _label(node, label, 10, INK)


def _draw_preview_visual(draw: ImageDraw.ImageDraw, item: Slide, font: ImageFont.ImageFont) -> None:
    labels = (item.visual_labels or item.content or [item.title])[:3]
    if item.visual_type in {"process", "timeline"}:
        for index, label in enumerate(labels):
            x = 755 + index * 172
            draw.polygon([(x, 285), (x + 135, 285), (x + 165, 355), (x + 135, 425), (x, 425), (x + 30, 355)], fill="#f4b21a" if index == 0 else "#17284a")
            draw.multiline_text((x + 38, 320), f"{index + 1}\n" + "\n".join(textwrap.wrap(label, 12)[:2]), fill="#0c1830" if index == 0 else "white", font=font, spacing=3)
    elif item.visual_type == "comparison":
        for index, label in enumerate(labels[:2]):
            x = 760 + index * 225; draw.rounded_rectangle((x, 230, x + 190, 535), radius=22, fill="#17284a" if index == 0 else "#f4b21a")
            draw.multiline_text((x + 20, 285), "\n".join(textwrap.wrap(label, 16)[:5]), fill="white" if index == 0 else "#0c1830", font=font, spacing=5)
    elif item.visual_type == "quote":
        draw.text((790, 175), "“", fill="#f4b21a", font=ImageFont.load_default(size=90)); draw.multiline_text((875, 330), "\n".join(textwrap.wrap(labels[0], 24)[:4]), fill="#0c1830", font=font, spacing=6)
    elif item.visual_type == "closing":
        draw.ellipse((850, 215, 1130, 495), fill="#fff7e0", outline="#f4b21a", width=7); draw.ellipse((930, 295, 1050, 415), fill="#f4b21a"); draw.text((970, 335), "→", fill="#0c1830", font=ImageFont.load_default(size=36))
    else:
        draw.line((825, 230, 975, 355), fill="#f4b21a", width=5); draw.line((1120, 230, 975, 355), fill="#f4b21a", width=5); draw.line((975, 355, 975, 520), fill="#f4b21a", width=5)
        draw.ellipse((890, 270, 1060, 440), fill="#17284a", outline="#f4b21a", width=5)
        draw.multiline_text((938, 325), "CORE\nIDEA", fill="#f4b21a", font=font, spacing=2)
        for index, label in enumerate(labels):
            x, y = ((750, 190), (1050, 190), (900, 485))[index]
            draw.ellipse((x, y, x + 145, y + 80), fill="#f4b21a")
            draw.multiline_text((x + 12, y + 18), "\n".join(textwrap.wrap(label, 14)[:2]), fill="#0c1830", font=font, spacing=3)


def render_deck(session_id: str, outline: list[Slide], root: Path) -> tuple[Path, list[Path]]:
    output = root / "storage" / session_id / "slides"; output.mkdir(parents=True, exist_ok=True)
    (output / "slide_content.json").write_text(json.dumps([s.model_dump() for s in outline], indent=2), encoding="utf-8")
    source_template = root.parent / "HR Transformation Mythbuster.pptx"
    if not source_template.exists():
        raise FileNotFoundError("HR Transformation Mythbuster.pptx is not available")
    path = output / "ekmind-content-deck.pptx"
    shutil.copy2(source_template, path)
    prs = Presentation(path)
    for slide_id in list(prs.slides._sldIdLst):
        relationship_id = slide_id.rId
        prs.part.drop_rel(relationship_id)
        prs.slides._sldIdLst.remove(slide_id)
    previews: list[Path] = []
    for item in outline:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        is_cover = item.slide_number == 1
        is_closing = item.slide_number == len(outline)
        if is_cover or is_closing:
            _background(slide, INK)
            gold_plane = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(6.85), Inches(0), Inches(3.15), Inches(5.625)); gold_plane.fill.solid(); gold_plane.fill.fore_color.rgb = GOLD; gold_plane.line.fill.background()
            _textbox(slide, "EKMIND  ·  MYTHBUSTER", 0.62, 0.48, 4.2, 0.3, 10, GOLD, True)
            _textbox(slide, item.title, 0.62, 1.22, 5.95, 2.0, 29 if is_cover else 31, WHITE, True)
            supporting = item.content[0] if item.content else ("Ideas that move from insight to action." if is_cover else "Turn the insight into one deliberate next move.")
            _textbox(slide, supporting, 0.65, 3.88, 5.25, 0.9, 15, WHITE)
            ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.35), Inches(1.55), Inches(1.9), Inches(1.9)); ring.fill.background(); ring.line.color.rgb = WHITE; ring.line.width = Pt(3)
            inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.9), Inches(2.1), Inches(0.8), Inches(0.8)); inner.fill.solid(); inner.fill.fore_color.rgb = INK; inner.line.fill.background()
        else:
            _background(slide, IVORY)
            field = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.3), Inches(0), Inches(4.7), Inches(5.625)); field.fill.solid(); field.fill.fore_color.rgb = PALE; field.line.fill.background()
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.52), Inches(1.48), Inches(0.07), Inches(3.38)); accent.fill.solid(); accent.fill.fore_color.rgb = GOLD; accent.line.fill.background()
            _textbox(slide, f"EKMIND  ·  {item.slide_number:02d}", 0.62, 0.36, 2.8, 0.25, 9, CORAL, True)
            _textbox(slide, item.title.upper(), 0.62, 0.72, 4.45, 0.82, 20, INK, True)
            y = 1.76
            for index, point in enumerate(item.content[:3], 1):
                _textbox(slide, str(index), 0.75, y, 0.35, 0.32, 13, GOLD, True)
                _textbox(slide, point, 1.18, y - 0.02, 3.72, 0.72, 15, INK)
                y += 0.96
            _textbox(slide, "MYTH  →  REALITY  →  ACTION", 0.62, 5.18, 3.3, 0.2, 8, CORAL, True)
            _add_visual(slide, item)

        notes = _clean_speaker_notes(item.speaker_notes)
        if not notes:
            notes = _humanized_fallback_notes(item.title, item.content, item.slide_number, len(outline))
        slide.notes_slide.notes_text_frame.text = notes

        image = Image.new("RGB", (1280, 720), "#0c1830" if is_cover or is_closing else "#faf8f2"); draw = ImageDraw.Draw(image)
        if is_cover or is_closing:
            draw.polygon([(880, 0), (1280, 0), (1280, 720), (980, 720)], fill="#f4b21a")
            draw.text((75, 62), "EKMIND  ·  MYTHBUSTER", fill="#f4b21a", font=ImageFont.load_default(size=13))
            draw.multiline_text((75, 160), "\n".join(textwrap.wrap(item.title, width=34)[:3]), fill="white", font=ImageFont.load_default(size=38), spacing=8)
            supporting = item.content[0] if item.content else "Ideas that move from insight to action."
            draw.multiline_text((78, 490), "\n".join(textwrap.wrap(supporting, width=48)[:2]), fill="white", font=ImageFont.load_default(size=18), spacing=5)
            draw.ellipse((955, 205, 1185, 435), outline="white", width=4); draw.ellipse((1025, 275, 1115, 365), fill="#0c1830")
        else:
            draw.rectangle((680, 0, 1280, 720), fill="#fff7e0")
            draw.rectangle((68, 160, 76, 620), fill="#f4b21a")
            draw.text((76, 42), f"EKMIND  ·  {item.slide_number:02d}", fill="#e4573d", font=ImageFont.load_default(size=12))
            draw.multiline_text((76, 82), "\n".join(textwrap.wrap(item.title.upper(), width=38)[:2]), fill="#0c1830", font=ImageFont.load_default(size=27), spacing=5)
            y = 220
            for index, point in enumerate(item.content[:3], 1):
                draw.text((95, y), str(index), fill="#f4b21a", font=ImageFont.load_default(size=14))
                draw.multiline_text((145, y), "\n".join(textwrap.wrap(point, width=36)[:3]), fill="#0c1830", font=ImageFont.load_default(size=19), spacing=6)
                y += 120
            _draw_preview_visual(draw, item, ImageFont.load_default(size=15))
        preview = output / f"slide-{item.slide_number:02d}.png"; image.save(preview); previews.append(preview)
    prs.save(path)
    if len(Presentation(path).slides) != len(outline):
        raise RuntimeError("PPTX validation failed")
    return path, previews
