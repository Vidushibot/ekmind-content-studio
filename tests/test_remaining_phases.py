from pathlib import Path
import subprocess

import imageio_ffmpeg
import pytest
from PIL import Image
from pptx import Presentation

from backend.evaluations.regression_runner import run_experiment
from backend.schemas.content import ApprovalRequest, GenerateRequest
from backend.schemas.content import Scene
from backend.services.content_service import ContentService
from backend.services.evaluation_service import human_edit_percentage, unsupported_claims
from backend.services.video_service import render_heygen_composite
from backend.providers.youtube import YouTubeProvider


def approved_session() -> tuple[ContentService, object]:
    service = ContentService(); view = service.create(GenerateRequest(topic="Responsible AI operating model", research_mode="Skip"))
    view = service.approve_angle(view.session_id, ApprovalRequest(angle_id=view.candidate_angles[0].angle_id))
    return service, service.approve_content(view.session_id, ApprovalRequest(edited_content=view.current_draft))


def test_factual_gate_and_human_diff():
    assert unsupported_claims("Research shows adoption improved 42%.", [])
    assert human_edit_percentage("abc", "abc") == 0


def test_all_governed_metrics_are_populated_without_invented_label_scores():
    service, view = approved_session()
    refreshed = service.get(view.session_id)
    assert len(refreshed.evaluation_metrics) == 24
    names = {row["metric"] for row in refreshed.evaluation_metrics}
    assert {"Retrieval Relevance", "Factual Correctness", "Human Edit %", "HITL Compliance", "Goal Completion"} <= names
    assert "Precision@K" not in names
    assert "Recall@K" not in names
    for metric_name in ("Evidence Coverage", "Groundedness", "Faithfulness", "Unsupported Claims"):
        metric = next(row for row in refreshed.evaluation_metrics if row["metric"] == metric_name)
        assert metric["value"] == "NOT_APPLICABLE"
        assert metric["status"] == "NOT_APPLICABLE"
        assert metric["target"] == "Research only"


def test_slide_and_video_hitl(tmp_path):
    service, view = approved_session()
    with pytest.raises(ValueError):
        service.generate_video(view.session_id)
    view = service.generate_slides(view.session_id)
    deck = Presentation(view.pptx_path)
    assert len(deck.slides) == 8
    assert all(slide.notes_slide.notes_text_frame.text.strip() for slide in deck.slides)
    assert all("http" not in slide.notes_slide.notes_text_frame.text for slide in deck.slides)
    assert len(view.slide_previews) == 8
    view = service.approve_slides(view.session_id)
    view = service.generate_video(view.session_id)
    assert Path(view.video_path).stat().st_size > 0
    assert view.video_validation["audio_stream"]
    assert view.video_validation["slide_assets_used"] == len(view.slide_outline)
    assert view.video_validation["composition"] == "approved_slide_previews"
    view = service.approve_video(view.session_id)
    assert view.status == "DONE"


def test_rejection_routes_are_recorded():
    service = ContentService(); view = service.create(GenerateRequest(topic="Responsible AI", research_mode="Skip"))
    view = service.reject_angles(view.session_id, ApprovalRequest(reason="Need sharper options"))
    assert view.status == "AWAITING_ANGLE_APPROVAL" and view.decision_history[-1]["decision"] == "REJECTED"
    view = service.approve_angle(view.session_id, ApprovalRequest(angle_id=view.candidate_angles[0].angle_id))
    view = service.reject_content(view.session_id, ApprovalRequest(reason="Try another angle"))
    assert view.status == "AWAITING_ANGLE_APPROVAL" and not view.angle_approved


def test_template_source_is_not_modified():
    template = Path(__file__).resolve().parents[2] / "HR Transformation Mythbuster.pptx"
    before = template.read_bytes()
    service, view = approved_session(); service.generate_slides(view.session_id)
    assert template.read_bytes() == before


def test_generated_deck_uses_mythbuster_reference_dimensions():
    service, view = approved_session(); view = service.generate_slides(view.session_id)
    deck = Presentation(view.pptx_path)
    assert round(deck.slide_width / 914400, 3) == 10.0
    assert round(deck.slide_height / 914400, 3) == 5.625


def test_golden_regression_runner(tmp_path):
    dataset = Path(__file__).resolve().parents[2] / "Ekmind_AI_Content_Studio_Golden_Dataset.xlsx"
    report = run_experiment(dataset, tmp_path, "baseline_v1")
    assert report["cases"] == 40
    assert (tmp_path / "baseline_v1.json").exists()


def test_heygen_compositor_uses_transitions_fullscreen_and_burned_captions(tmp_path):
    slides = []
    for index, color in enumerate(("#172554", "#0f766e", "#7c2d12"), 1):
        path = tmp_path / f"slide-{index}.png"
        Image.new("RGB", (1280, 720), color).save(path)
        slides.append(str(path))
    avatar = tmp_path / "avatar.mp4"
    subprocess.run([imageio_ffmpeg.get_ffmpeg_exe(), "-y", "-f", "lavfi", "-i", "testsrc2=size=1280x720:rate=30", "-f", "lavfi", "-i", "sine=frequency=440:sample_rate=48000", "-t", "9", "-c:v", "libx264", "-c:a", "aac", str(avatar)], check=True, capture_output=True, timeout=120)
    scenes = [Scene(scene_type="SLIDE_FULL", spoken_text=f"Caption for scene {index} with enough words for timing.", slide_number=index) for index in range(1, 4)]
    video, srt, validation = render_heygen_composite("composition-test", scenes, slides, avatar, 9, tmp_path)
    assert video.stat().st_size > 0 and srt.exists()
    assert validation["slide_transitions"] == "crossfade"
    assert validation["captions"] == "burned_in_and_srt"
    assert validation["composition"] == "full_screen_intro_outro_with_large_avatar_pip"


def test_youtube_upload_is_always_private(tmp_path, monkeypatch):
    captured = {}

    class UploadRequest:
        def next_chunk(self):
            return None, {"id": "private-video-id"}

    class Videos:
        def insert(self, **kwargs):
            captured.update(kwargs)
            return UploadRequest()

    class Service:
        def videos(self):
            return Videos()

    video = tmp_path / "approved.mp4"
    video.write_bytes(b"video")
    provider = YouTubeProvider("client", "secret", "refresh")
    monkeypatch.setattr(provider, "_service", lambda: Service())
    result = provider.upload_private(video, "Approved title", "Approved description")

    assert captured["body"]["status"]["privacyStatus"] == "private"
    assert result["privacy_status"] == "private"
    assert result["video_id"] == "private-video-id"
